"""
Generate a styled PowerPoint slide with the 10-trial random search table.

Usage:
    pip install python-pptx
    python make_table_pptx.py

Output:
    table_10trials_combined.pptx  (saved next to this script)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ---------------------------------------------------------------
# Data
# ---------------------------------------------------------------
HEADERS = ["Trial", "K", "hidden", "dropout", "lr", "batch",
           "val NLL", "val MSE", "test NLL", "test MSE"]

ROWS = [
    ["0",   "2", "64",  "0.2", "5e-4", "256", "1.1984", "0.5499", "2.1549", "1.2526"],
    ["1",   "2", "64",  "0.3", "5e-4", "256", "1.2941", "0.5495", "2.2971", "1.2499"],
    ["2 ★", "3", "64",  "0.2", "1e-3", "256", "0.9092", "0.5498", "2.0379", "1.2523"],
    ["3",   "3", "128", "0.2", "1e-3", "512", "1.0661", "0.5509", "1.9908", "1.2577"],
    ["4",   "5", "128", "0.3", "5e-4", "512", "1.3547", "0.5505", "2.2224", "1.2516"],
    ["5",   "5", "128", "0.3", "1e-3", "512", "1.1449", "0.5488", "1.9922", "1.2535"],
    ["6",   "5", "64",  "0.2", "5e-4", "256", "1.0583", "0.5488", "1.9042", "1.2531"],
    ["7",   "5", "128", "0.3", "5e-4", "512", "1.3547", "0.5505", "2.2224", "1.2516"],
    ["8",   "5", "64",  "0.2", "1e-3", "512", "0.9368", "0.5501", "1.7532", "1.2524"],
    ["9",   "2", "128", "0.2", "5e-4", "256", "1.1659", "0.5513", "2.2296", "1.2537"],
]
BEST_ROW_INDEX = 2  # Trial 2

# Colors
NAVY        = RGBColor(0x1E, 0x27, 0x61)
NAVY_LIGHT  = RGBColor(0x3C, 0x4F, 0x7F)
GREEN_HEAD  = RGBColor(0x0F, 0x6E, 0x56)
RED_HEAD    = RGBColor(0x99, 0x3C, 0x1D)
BEST_FILL   = RGBColor(0xE1, 0xF5, 0xEE)
BEST_TEXT   = RGBColor(0x0F, 0x6E, 0x56)
ZEBRA_GRAY  = RGBColor(0xF7, 0xF7, 0xF7)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x1A, 0x1A, 0x1A)
ICE_BLUE    = RGBColor(0xCA, 0xDC, 0xFC)
BORDER_GRAY = RGBColor(0xDD, 0xDD, 0xDD)


def set_cell_fill(cell, rgb):
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb


def style_run(run, *, size=12, bold=False, color=BLACK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def write_cell(cell, text, *, size=12, bold=False, color=BLACK,
               fill=None, align=PP_ALIGN.CENTER):
    if fill is not None:
        set_cell_fill(cell, fill)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.05)
    cell.margin_right = Inches(0.05)
    cell.margin_top = Inches(0.03)
    cell.margin_bottom = Inches(0.03)

    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    run = p.add_run()
    run.text = text
    style_run(run, size=size, bold=bold, color=color)


def add_text_box(slide, *, x, y, w, h, text, size=14, bold=False,
                 color=BLACK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    style_run(run, size=size, bold=bold, color=color, font=font)
    return box


def add_filled_rect(slide, *, x, y, w, h, fill, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def main():
    out_dir = Path(__file__).parent.resolve()
    out_path = out_dir / "table_10trials_combined.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Header bar
    add_filled_rect(slide, x=Inches(0), y=Inches(0),
                    w=prs.slide_width, h=Inches(0.85), fill=NAVY)
    add_text_box(slide, x=Inches(0.5), y=Inches(0.18), w=Inches(10), h=Inches(0.55),
                 text="Random Search Results — 10 Trials (NLL & MSE)",
                 size=24, bold=True, color=WHITE)
    add_text_box(slide, x=Inches(8.5), y=Inches(0.28), w=Inches(4.5), h=Inches(0.4),
                 text="MDN Surrogate Model | 754 ion-collision events",
                 size=12, color=ICE_BLUE, align=PP_ALIGN.RIGHT)

    # Table
    n_rows = len(ROWS) + 1
    n_cols = len(HEADERS)
    table_left = Inches(0.5)
    table_top = Inches(1.15)
    table_width = Inches(12.333)
    table_height = Inches(4.6)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, table_left, table_top, table_width, table_height
    )
    table = table_shape.table

    # Column widths (sum = 12.333")
    col_widths_in = [0.85, 0.7, 1.0, 1.0, 1.0, 1.0, 1.6, 1.6, 1.6, 1.6]
    assert abs(sum(col_widths_in) - 12.95) < 0.5
    for i, w in enumerate(col_widths_in):
        table.columns[i].width = Inches(w)

    # Row heights
    table.rows[0].height = Inches(0.5)
    body_h = Inches(0.4)
    for r in range(1, n_rows):
        table.rows[r].height = body_h

    # Header row
    for c, header in enumerate(HEADERS):
        cell = table.cell(0, c)
        if c == 0:
            fill = NAVY
        elif c <= 5:
            fill = NAVY_LIGHT
        elif c <= 7:
            fill = GREEN_HEAD
        else:
            fill = RED_HEAD
        write_cell(cell, header, size=14, bold=True, color=WHITE, fill=fill)

    # Body rows
    for r, row in enumerate(ROWS, start=1):
        is_best = (r - 1 == BEST_ROW_INDEX)
        if is_best:
            row_fill = BEST_FILL
            text_color = BEST_TEXT
            bold = True
        else:
            row_fill = ZEBRA_GRAY if (r % 2 == 1) else WHITE
            text_color = BLACK
            bold = False

        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell_bold = bold or (c == 0)  # first column always slightly bold
            write_cell(cell, val, size=13, bold=cell_bold,
                       color=text_color, fill=row_fill)

    # Best-config callout
    callout_top = Inches(6.0)
    add_filled_rect(slide,
                    x=Inches(0.5), y=callout_top,
                    w=Inches(12.333), h=Inches(1.25),
                    fill=BEST_FILL, line_color=GREEN_HEAD)

    add_text_box(slide,
                 x=Inches(0.7), y=callout_top + Inches(0.1),
                 w=Inches(12), h=Inches(0.4),
                 text="★ Best Configuration (Trial 2)",
                 size=16, bold=True, color=BEST_TEXT)

    add_text_box(slide,
                 x=Inches(0.7), y=callout_top + Inches(0.45),
                 w=Inches(12), h=Inches(0.35),
                 text="K=3  |  hidden=64  |  dropout=0.2  |  lr=1e-3  |  batch=256   →   val NLL = 0.909  |  test NLL = 2.04  |  test MSE = 1.25",
                 size=13, bold=True, color=BEST_TEXT)

    add_text_box(slide,
                 x=Inches(0.7), y=callout_top + Inches(0.78),
                 w=Inches(12), h=Inches(0.4),
                 text=("Selection criterion: validation NLL. "
                       "MSE is nearly identical across trials (mixture means converge similarly) — "
                       "NLL captures distribution shape (σ, π, multimodality), which is what differs."),
                 size=11, color=BLACK)

    prs.save(str(out_path))
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
